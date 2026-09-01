"""
ppt/parse_ppt.py
Usage:  python parse_ppt.py <filename>
Place file in input/ folder. Output goes to output/ with a datetime stamp.
"""

import sys, os
from datetime import datetime

def main():
    if len(sys.argv) < 2:
        print("Usage: python parse_ppt.py <filename>")
        print("Example: python parse_ppt.py deck.pptx")
        sys.exit(1)

    filename  = sys.argv[1].strip()
    script_dir = os.path.dirname(os.path.abspath(__file__))
    input_path = os.path.join(script_dir, "input", filename)

    if not os.path.isfile(input_path):
        print(f"ERROR: File not found: {input_path}")
        print(f"Place '{filename}' inside the input/ folder and try again.")
        sys.exit(1)

    # ── install dependency if missing ──────────────────────────────────────
    try:
        from pptx import Presentation
    except ImportError:
        print("python-pptx not found. Installing...")
        os.system("pip install python-pptx --break-system-packages -q")
        from pptx import Presentation

    lines = []

    # ── open file ──────────────────────────────────────────────────────────
    try:
        prs = Presentation(input_path)
    except Exception as e:
        print(f"ERROR opening file: {e}")
        sys.exit(1)

    lines.append(f"FILE: {filename}")
    lines.append(f"PARSED AT: {datetime.now()}")

    # ── core properties ────────────────────────────────────────────────────
    lines.append("\n" + "="*60)
    lines.append("CORE PROPERTIES")
    lines.append("="*60)
    try:
        cp = prs.core_properties
        for attr in ["title","author","subject","keywords","description",
                     "category","last_modified_by","created","modified",
                     "revision"]:
            try:
                val = getattr(cp, attr, None)
                lines.append(f"{attr}: {val}")
            except Exception as e:
                lines.append(f"{attr}: [ERROR: {e}]")
    except Exception as e:
        lines.append(f"[Could not read core properties: {e}]")

    # ── slide dimensions ───────────────────────────────────────────────────
    try:
        lines.append(f"\nslide_width_emu: {prs.slide_width}")
        lines.append(f"slide_height_emu: {prs.slide_height}")
    except Exception as e:
        lines.append(f"[Could not read slide dimensions: {e}]")

    # ── slides ─────────────────────────────────────────────────────────────
    try:
        total = len(prs.slides)
    except Exception as e:
        total = "unknown"
    lines.append(f"\ntotal_slides: {total}")

    for slide_idx, slide in enumerate(prs.slides):
        lines.append("\n" + "="*60)
        lines.append(f"SLIDE {slide_idx + 1}")
        lines.append("="*60)

        # layout name
        try:
            lines.append(f"layout: {slide.slide_layout.name}")
        except Exception as e:
            lines.append(f"layout: [ERROR: {e}]")

        # slide id
        try:
            lines.append(f"slide_id: {slide.slide_id}")
        except Exception as e:
            lines.append(f"slide_id: [ERROR: {e}]")

        # shapes
        try:
            shape_list = list(slide.shapes)
        except Exception as e:
            lines.append(f"[Could not iterate shapes: {e}]")
            shape_list = []

        for shape_idx, shape in enumerate(shape_list):
            lines.append(f"\n  -- Shape {shape_idx + 1} --")

            try: lines.append(f"  name: {shape.name}")
            except Exception as e: lines.append(f"  name: [ERROR: {e}]")

            try: lines.append(f"  shape_type: {shape.shape_type}")
            except Exception as e: lines.append(f"  shape_type: [ERROR: {e}]")

            try: lines.append(f"  left: {shape.left}  top: {shape.top}  width: {shape.width}  height: {shape.height}")
            except Exception as e: lines.append(f"  position/size: [ERROR: {e}]")

            # text frame — dump every run of every paragraph, no filtering
            try:
                if shape.has_text_frame:
                    lines.append("  [TEXT FRAME]")
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        try:
                            full_para_text = para.text   # entire paragraph at once
                            lines.append(f"    para[{para_idx}]: {full_para_text}")
                            # also dump individual runs for completeness
                            for run_idx, run in enumerate(para.runs):
                                try:
                                    lines.append(f"      run[{run_idx}]: {run.text}")
                                except Exception as e:
                                    lines.append(f"      run[{run_idx}]: [ERROR: {e}]")
                        except Exception as e:
                            lines.append(f"    para[{para_idx}]: [ERROR: {e}]")
            except Exception as e:
                lines.append(f"  text_frame: [ERROR: {e}]")

            # table — dump every cell
            try:
                if shape.has_table:
                    lines.append("  [TABLE]")
                    for row_idx, row in enumerate(shape.table.rows):
                        row_cells = []
                        for cell in row.cells:
                            try:
                                row_cells.append(cell.text)
                            except Exception as e:
                                row_cells.append(f"[ERROR: {e}]")
                        lines.append(f"    row[{row_idx}]: {' | '.join(row_cells)}")
            except Exception as e:
                lines.append(f"  table: [ERROR: {e}]")

            # chart
            try:
                if shape.has_chart:
                    chart = shape.chart
                    lines.append("  [CHART]")
                    try:
                        lines.append(f"    chart_type: {chart.chart_type}")
                    except Exception as e:
                        lines.append(f"    chart_type: [ERROR: {e}]")
                    try:
                        title_text = chart.chart_title.text_frame.text if chart.has_title else "(no title)"
                        lines.append(f"    chart_title: {title_text}")
                    except Exception as e:
                        lines.append(f"    chart_title: [ERROR: {e}]")
                    # series names and values
                    try:
                        for series in chart.series:
                            try:
                                lines.append(f"    series: {series.name}")
                            except Exception as e:
                                lines.append(f"    series: [ERROR: {e}]")
                    except Exception as e:
                        lines.append(f"    series: [ERROR: {e}]")
            except Exception as e:
                lines.append(f"  chart: [ERROR: {e}]")

        # speaker notes
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                lines.append(f"\n  [SPEAKER NOTES]\n  {notes_text}")
        except Exception as e:
            lines.append(f"  notes: [ERROR: {e}]")

    # ── write output ───────────────────────────────────────────────────────
    timestamp  = datetime.now().strftime("%Y%m%d_%H%M%S")
    base       = os.path.splitext(filename)[0]
    output_dir = os.path.join(script_dir, "output")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, f"{base}__{timestamp}.txt")

    try:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))
        print(f"Done! Output → {output_path}")
    except Exception as e:
        print(f"ERROR writing output: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
